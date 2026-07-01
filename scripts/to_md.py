#!/usr/bin/env python3
"""
텍스트 → Markdown 변환기 (이미지 제외)

지원 형식:
  로컬 파일: PDF, DOCX, PPTX, XLSX, TXT, HTML, RTF
  Google Drive: Google Docs, Google Sheets, Google Slides, 업로드된 PDF/DOCX 등

사용법:
  # 로컬 파일
  python to_md.py 파일.pdf
  python to_md.py 폴더/
  python to_md.py 파일.pdf -o obsidian/20-원자료/인박스/

  # Google Drive (파일 ID 또는 공유 URL)
  python to_md.py "https://docs.google.com/document/d/FILE_ID/edit"
  python to_md.py --drive FILE_ID -o obsidian/20-원자료/인박스/

  # Google Drive 폴더 전체
  python to_md.py --drive-folder FOLDER_ID -o obsidian/20-원자료/인박스/

Google Drive 최초 사용:
  pip install google-api-python-client google-auth-httplib2 google-auth-oauthlib
  → credentials.json 을 이 스크립트와 같은 폴더에 놓고 실행
  → 브라우저 인증 1회 후 token.json 자동 저장
"""

import sys
import re
import io
from pathlib import Path
from datetime import date


# ─── 텍스트 클리닝 ────────────────────────────────────────────────────────────

def clean(text: str) -> str:
    text = re.sub(r'\r\n', '\n', text)
    text = re.sub(r'\n{3,}', '\n\n', text)
    text = "\n".join(line.rstrip() for line in text.splitlines())
    return text.strip()


def wrap_md(title: str, text: str) -> str:
    today = date.today().isoformat()
    return f"""---
date: {today}
source: import
tags: [원자료]
---

# {title}

{text}
"""


# ─── 로컬 파일 추출 ───────────────────────────────────────────────────────────

def from_pdf(path: Path) -> str:
    try:
        import fitz
        doc = fitz.open(path)
        return "\n\n".join(p.get_text("text").strip() for p in doc if p.get_text("text").strip())
    except ImportError:
        pass
    try:
        import pdfplumber
        with pdfplumber.open(path) as pdf:
            return "\n\n".join(
                p.extract_text().strip() for p in pdf.pages if p.extract_text()
            )
    except ImportError:
        sys.exit("pip install pymupdf  또는  pip install pdfplumber")


def from_docx(path: Path) -> str:
    try:
        from docx import Document
    except ImportError:
        sys.exit("pip install python-docx")
    doc = Document(path)
    return "\n\n".join(p.text.strip() for p in doc.paragraphs if p.text.strip())


def from_pptx(path: Path) -> str:
    try:
        from pptx import Presentation
    except ImportError:
        sys.exit("pip install python-pptx")
    prs = Presentation(path)
    slides = []
    for i, slide in enumerate(prs.slides, 1):
        texts = [shape.text.strip() for shape in slide.shapes if hasattr(shape, "text") and shape.text.strip()]
        if texts:
            slides.append(f"## 슬라이드 {i}\n\n" + "\n\n".join(texts))
    return "\n\n".join(slides)


def from_xlsx(path: Path) -> str:
    if path.suffix.lower() == ".xls":
        try:
            import xlrd
            wb = xlrd.open_workbook(path)
            sheets = []
            for sheet in wb.sheets():
                rows = []
                for i in range(sheet.nrows):
                    cells = [str(sheet.cell_value(i, j)) for j in range(sheet.ncols)
                             if str(sheet.cell_value(i, j)).strip()]
                    if cells:
                        rows.append("  ".join(cells))
                if rows:
                    sheets.append(f"## {sheet.name}\n\n" + "\n".join(rows))
            return "\n\n".join(sheets)
        except ImportError:
            sys.exit("pip3 install --break-system-packages xlrd")

    try:
        import openpyxl
    except ImportError:
        sys.exit("pip install openpyxl")
    wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
    sheets = []
    for sheet in wb.worksheets:
        rows = []
        for row in sheet.iter_rows(values_only=True):
            cells = [str(c) for c in row if c is not None]
            if cells:
                rows.append("  ".join(cells))
        if rows:
            sheets.append(f"## {sheet.title}\n\n" + "\n".join(rows))
    return "\n\n".join(sheets)


def from_txt(path: Path) -> str:
    for enc in ("utf-8", "euc-kr", "cp949"):
        try:
            return path.read_text(encoding=enc)
        except (UnicodeDecodeError, LookupError):
            continue
    return path.read_text(errors="ignore")


def from_html(path: Path) -> str:
    try:
        from bs4 import BeautifulSoup
    except ImportError:
        sys.exit("pip install beautifulsoup4")
    soup = BeautifulSoup(path.read_text(encoding="utf-8", errors="ignore"), "html.parser")
    for tag in soup(["script", "style", "img"]):
        tag.decompose()
    return soup.get_text(separator="\n")


def from_hwpx(path: Path) -> str:
    import zipfile
    import xml.etree.ElementTree as ET

    texts = []
    with zipfile.ZipFile(path, "r") as z:
        # 섹션 파일 정렬 추출 (section0.xml, section1.xml ...)
        section_files = sorted(
            [f for f in z.namelist() if re.match(r"Contents/section\d+\.xml", f)]
        )
        if not section_files:
            # 일부 HWPX는 경로가 다를 수 있음
            section_files = sorted(
                [f for f in z.namelist() if f.endswith(".xml") and "section" in f.lower()]
            )
        for sf in section_files:
            xml_bytes = z.read(sf)
            root = ET.fromstring(xml_bytes)
            # 한글 XML 네임스페이스 내 hp:t 태그에서 텍스트 추출
            for elem in root.iter():
                if elem.tag.endswith("}t") or elem.tag == "t":
                    if elem.text and elem.text.strip():
                        texts.append(elem.text.strip())
    return "\n\n".join(texts)


EXTRACTORS = {
    ".pdf":  from_pdf,
    ".docx": from_docx,
    ".doc":  from_docx,
    ".pptx": from_pptx,
    ".ppt":  from_pptx,
    ".xlsx": from_xlsx,
    ".xls":  from_xlsx,
    ".txt":  from_txt,
    ".md":   from_txt,
    ".html": from_html,
    ".htm":  from_html,
    ".hwpx": from_hwpx,
    ".hwpz": from_hwpx,
}


SKIPPED_LOG = Path(__file__).parent / ".skipped.txt"


def log_skipped(path: Path):
    with open(SKIPPED_LOG, "a", encoding="utf-8") as f:
        f.write(str(path) + "\n")


def convert_local(path: Path, output_dir: Path, timeout: int = 60, skip_existing: bool = True):
    import threading

    ext = path.suffix.lower()
    extractor = EXTRACTORS.get(ext)
    if not extractor:
        print(f"  건너뜀 (미지원 형식): {path.name}")
        return

    out = output_dir / (path.stem + ".md")
    if skip_existing and out.exists():
        print(f"  건너뜀 (이미 변환됨): {path.name}")
        return

    print(f"변환 중: {path.name}")

    result = [None]
    error  = [None]

    def _run():
        try:
            result[0] = clean(extractor(path))
        except Exception as e:
            error[0] = e

    t = threading.Thread(target=_run, daemon=True)
    t.start()
    t.join(timeout)

    if t.is_alive():
        print(f"  건너뜀 (60초 초과): {path.name}")
        log_skipped(path)
        return
    if error[0]:
        print(f"  건너뜀 (오류): {path.name} — {error[0]}")
        log_skipped(path)
        return

    md = wrap_md(path.stem, result[0])
    output_dir.mkdir(parents=True, exist_ok=True)
    out.write_text(md, encoding="utf-8")
    print(f"  저장: {out}")


# ─── Google Drive ─────────────────────────────────────────────────────────────

GDRIVE_EXPORT_MIME = {
    "application/vnd.google-apps.document":     ("text/plain", ".txt"),
    "application/vnd.google-apps.spreadsheet":  ("text/csv",   ".csv"),
    "application/vnd.google-apps.presentation": ("text/plain", ".txt"),
}

GDRIVE_DOWNLOAD_MIME = {
    "application/pdf": ".pdf",
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document": ".docx",
    "application/vnd.openxmlformats-officedocument.presentationml.presentation": ".pptx",
    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet": ".xlsx",
    "text/plain": ".txt",
    "text/html": ".html",
}

SCOPES = ["https://www.googleapis.com/auth/drive.readonly"]


def get_drive_service():
    from google.oauth2.credentials import Credentials
    from google_auth_oauthlib.flow import InstalledAppFlow
    from google.auth.transport.requests import Request
    from googleapiclient.discovery import build

    creds = None
    token_path = Path(__file__).parent / "token.json"
    creds_path = Path(__file__).parent / "credentials.json"

    if token_path.exists():
        creds = Credentials.from_authorized_user_file(token_path, SCOPES)
    if not creds or not creds.valid:
        if creds and creds.expired and creds.refresh_token:
            creds.refresh(Request())
        else:
            if not creds_path.exists():
                sys.exit(
                    "credentials.json 없음.\n"
                    "Google Cloud Console → API 및 서비스 → 사용자 인증 정보 → OAuth 2.0 클라이언트 ID\n"
                    "→ '데스크톱 앱' 유형으로 생성 후 credentials.json 을 이 스크립트 폴더에 저장"
                )
            flow = InstalledAppFlow.from_client_secrets_file(creds_path, SCOPES)
            creds = flow.run_local_server(port=0)
        token_path.write_text(creds.to_json())
    return build("drive", "v3", credentials=creds)


def extract_file_id(url_or_id: str) -> str:
    patterns = [
        r"/d/([a-zA-Z0-9_-]{25,})",
        r"id=([a-zA-Z0-9_-]{25,})",
        r"^([a-zA-Z0-9_-]{25,})$",
    ]
    for p in patterns:
        m = re.search(p, url_or_id)
        if m:
            return m.group(1)
    sys.exit(f"파일 ID를 파싱할 수 없습니다: {url_or_id}")


def convert_drive_file(service, file_id: str, output_dir: Path):
    from googleapiclient.http import MediaIoBaseDownload

    meta = service.files().get(fileId=file_id, fields="id,name,mimeType").execute()
    name = meta["name"]
    mime = meta["mimeType"]
    print(f"Drive 변환 중: {name} ({mime})")

    if mime in GDRIVE_EXPORT_MIME:
        export_mime, ext = GDRIVE_EXPORT_MIME[mime]
        data = service.files().export(fileId=file_id, mimeType=export_mime).execute()
        tmp = Path(f"/tmp/{name}{ext}")
        tmp.write_bytes(data if isinstance(data, bytes) else data.encode("utf-8"))
    elif mime in GDRIVE_DOWNLOAD_MIME:
        ext = GDRIVE_DOWNLOAD_MIME[mime]
        req = service.files().get_media(fileId=file_id)
        buf = io.BytesIO()
        dl = MediaIoBaseDownload(buf, req)
        done = False
        while not done:
            _, done = dl.next_chunk()
        tmp = Path(f"/tmp/{name}{ext}")
        tmp.write_bytes(buf.getvalue())
    else:
        print(f"  건너뜀 (미지원 MIME): {mime}")
        return

    convert_local(tmp, output_dir)
    tmp.unlink(missing_ok=True)


def convert_drive_folder(service, folder_id: str, output_dir: Path,
                         history: set = None, recursive: bool = True):
    """폴더 내 파일 변환. recursive=True면 하위 폴더도 처리."""
    if history is None:
        history = set()

    supported_mimes = set(GDRIVE_EXPORT_MIME) | set(GDRIVE_DOWNLOAD_MIME)
    query = f"'{folder_id}' in parents and trashed=false"
    page_token = None

    while True:
        resp = service.files().list(
            q=query,
            fields="nextPageToken, files(id,name,mimeType)",
            pageSize=200,
            pageToken=page_token,
        ).execute()

        for f in resp.get("files", []):
            if f["mimeType"] == "application/vnd.google-apps.folder":
                if recursive:
                    sub_dir = output_dir / _safe_name(f["name"])
                    convert_drive_folder(service, f["id"], sub_dir, history, recursive)
            elif f["mimeType"] in supported_mimes:
                if f["id"] in history:
                    print(f"  건너뜀 (이미 변환됨): {f['name']}")
                    continue
                try:
                    convert_drive_file(service, f["id"], output_dir)
                    history.add(f["id"])
                except Exception as e:
                    print(f"  오류 ({f['name']}): {e}")

        page_token = resp.get("nextPageToken")
        if not page_token:
            break


def convert_drive_all(service, output_dir: Path):
    """Google Drive 전체 파일 변환 (이력 관리로 중복 방지)."""
    import json

    history_file = Path(__file__).parent / ".drive_history.json"
    history: set = set()
    if history_file.exists():
        try:
            history = set(json.loads(history_file.read_text()))
        except Exception:
            pass

    supported_mimes = set(GDRIVE_EXPORT_MIME) | set(GDRIVE_DOWNLOAD_MIME)
    mime_filter = " or ".join(f"mimeType='{m}'" for m in supported_mimes)
    query = f"trashed=false and ({mime_filter})"

    page_token = None
    total = 0
    skipped = 0

    print("Google Drive 전체 파일 목록 수집 중...")

    while True:
        resp = service.files().list(
            q=query,
            fields="nextPageToken, files(id,name,mimeType,parents)",
            pageSize=200,
            pageToken=page_token,
        ).execute()

        files = resp.get("files", [])
        for f in files:
            total += 1
            if f["id"] in history:
                skipped += 1
                continue
            try:
                convert_drive_file(service, f["id"], output_dir)
                history.add(f["id"])
            except Exception as e:
                print(f"  오류 ({f['name']}): {e}")

        # 중간 저장 (중단돼도 이력 유지)
        history_file.write_text(json.dumps(list(history)))

        page_token = resp.get("nextPageToken")
        if not page_token:
            break

    print(f"\n전체: {total}개 | 신규 변환: {total - skipped}개 | 건너뜀: {skipped}개")
    history_file.write_text(json.dumps(list(history)))


def _safe_name(name: str) -> str:
    return re.sub(r'[<>:"/\\|?*]', '_', name)


# ─── 메인 ─────────────────────────────────────────────────────────────────────

def main():
    import argparse
    parser = argparse.ArgumentParser(description="텍스트 파일 → Markdown 변환기")
    parser.add_argument("input", nargs="?", help="로컬 파일/폴더 또는 Google Drive URL/ID")
    parser.add_argument("-o", "--output", default=None, help="출력 폴더")
    parser.add_argument("--drive", metavar="ID_OR_URL", help="Google Drive 파일 ID 또는 URL")
    parser.add_argument("--drive-folder", metavar="FOLDER_ID", help="Google Drive 특정 폴더 ID")
    parser.add_argument("--drive-all", action="store_true", help="Google Drive 전체 파일 변환")
    parser.add_argument("--retry", action="store_true", help="이전에 건너뛴 파일 재시도 (타임아웃 2배)")
    args = parser.parse_args()

    output_dir = Path(args.output) if args.output else None

    # 건너뛴 파일 재시도 모드
    if args.retry:
        if not SKIPPED_LOG.exists():
            print("건너뛴 파일 기록이 없습니다.")
            return
        files = [Path(p.strip()) for p in SKIPPED_LOG.read_text().splitlines() if p.strip()]
        if not files:
            print("재시도할 파일이 없습니다.")
            return
        out = output_dir or Path(".")
        print(f"재시도: {len(files)}개 파일 (타임아웃 120초)")
        SKIPPED_LOG.unlink()
        for f in files:
            if f.exists():
                convert_local(f, out, timeout=120, skip_existing=False)
            else:
                print(f"  파일 없음: {f}")
        print("완료.")
        return

    # Google Drive 전체 모드
    if args.drive_all:
        service = get_drive_service()
        out = output_dir or Path(".")
        convert_drive_all(service, out)
        print("완료.")
        return

    # Google Drive 단일/폴더 모드
    if args.drive or args.drive_folder:
        service = get_drive_service()
        out = output_dir or Path(".")
        if args.drive:
            fid = extract_file_id(args.drive)
            convert_drive_file(service, fid, out)
        else:
            convert_drive_folder(service, args.drive_folder, out)
        print("완료.")
        return

    # Google Drive URL이 input으로 들어온 경우
    if args.input and ("docs.google.com" in args.input or "drive.google.com" in args.input):
        service = get_drive_service()
        fid = extract_file_id(args.input)
        out = output_dir or Path(".")
        convert_drive_file(service, fid, out)
        print("완료.")
        return

    # 로컬 파일/폴더 모드
    if not args.input:
        parser.print_help()
        sys.exit(1)

    input_path = Path(args.input)
    if input_path.is_dir():
        files = [f for f in input_path.rglob("*") if f.suffix.lower() in EXTRACTORS]
        if not files:
            sys.exit("변환할 파일이 없습니다.")
        for f in files:
            out = output_dir or f.parent
            convert_local(f, out)
    elif input_path.is_file():
        out = output_dir or input_path.parent
        convert_local(input_path, out)
    else:
        sys.exit(f"파일을 찾을 수 없습니다: {input_path}")

    print("완료.")


if __name__ == "__main__":
    main()
